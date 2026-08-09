#!/usr/bin/env python3
"""把镜像里的 PPT 课件转成门户专属的 Slides JSON。

用法:
  convert-henukit-slides.py --mirror <public_root> --out <slides_dir> \
    --manifest <manifest.json>

对 manifest 中 role 为 "课件PPT" 的资产:
  - .pptx 直接用 python-pptx 按页抽取文本;
  - .ppt 先用 LibreOffice headless 转成 .pptx 再抽取(需要 soffice)。

输出 <slides_dir>/<storage_key>.json:
  {"slides": [{"title": str, "blocks": [str, ...]}, ...]}

依赖: python3-pptx;处理 .ppt 还需要 libreoffice-impress。
依赖缺失时按退出码 2/3 结束；发布激活必须失败关闭，不能跳过派生产物;
镜像文件本身永远不会被修改。

幂等:输出已存在且不早于源文件时跳过,避免每次同步都重转。
"""

import argparse
import json
import pathlib
import shutil
import subprocess
import sys

REVIEW_MARKER = "待复核"
SLIDE_ROLE = "课件PPT"


def fail(message: str, code: int) -> None:
    print(f"convert-henukit-slides: {message}", file=sys.stderr)
    sys.exit(code)


def extract_pptx(path: pathlib.Path) -> list[dict]:
    try:
        from pptx import Presentation
    except ImportError:
        fail(
            "python3-pptx is required to extract slide text "
            "(install python3-pptx, e.g. apt install python3-pptx)",
            2,
        )

    presentation = Presentation(str(path))
    slides = []
    for slide in presentation.slides:
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        texts.append(text)
            elif getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    cells = [
                        cell.text.strip()
                        for cell in row.cells
                        if cell.text.strip()
                    ]
                    if cells:
                        texts.append(" | ".join(cells))
        # 去掉连续重复段(常见于文本框与备注复制)。
        deduped: list[str] = []
        for text in texts:
            if not deduped or deduped[-1] != text:
                deduped.append(text)
        if deduped:
            slides.append({"title": deduped[0], "blocks": deduped[1:]})
    return slides


def convert_legacy_ppt(source: pathlib.Path, scratch: pathlib.Path) -> pathlib.Path:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice is None:
        fail(
            f"LibreOffice (soffice) is required to convert legacy .ppt files "
            f"(needed for {source.name}); install libreoffice-impress",
            3,
        )
    scratch.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        [soffice, "--headless", "--convert-to", "pptx", "--outdir", str(scratch), str(source)],
        check=True,
        capture_output=True,
        timeout=300,
    )
    converted = scratch / f"{source.stem}.pptx"
    if not converted.is_file():
        fail(f"soffice did not produce {converted.name}", 4)
    return converted


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--mirror", required=True, help="镜像根目录(nginx 服务的目录)")
    parser.add_argument("--out", required=True, help="幻灯片 JSON 输出目录")
    parser.add_argument("--manifest", required=True, help="HENU-Final-Review manifest.json")
    args = parser.parse_args()

    mirror = pathlib.Path(args.mirror).resolve()
    out = pathlib.Path(args.out).resolve()
    manifest = pathlib.Path(args.manifest).resolve()

    if not manifest.is_file():
        fail(f"manifest not found: {manifest}", 4)
    data = json.loads(manifest.read_text(encoding="utf-8"))

    scratch = out / ".conversion-tmp"
    converted = 0
    skipped = 0
    failed = []

    for subject in data.get("subjects", []):
        for asset in subject.get("assets", []):
            role = asset.get("role", "")
            if role != SLIDE_ROLE or role.startswith(REVIEW_MARKER):
                continue
            public_path = asset.get("publicPath", "")
            if not public_path:
                continue
            source = (mirror / public_path).resolve()
            if not source.is_relative_to(mirror) or not source.is_file():
                failed.append(f"{public_path}: source missing")
                continue

            target = out / f"{public_path}.json"
            if target.is_file() and target.stat().st_mtime >= source.stat().st_mtime:
                skipped += 1
                continue

            target.parent.mkdir(parents=True, exist_ok=True)
            try:
                work = source
                if source.suffix.lower() == ".ppt":
                    work = convert_legacy_ppt(source, scratch)
                slides = extract_pptx(work)
                target.write_text(
                    json.dumps({"slides": slides}, ensure_ascii=False),
                    encoding="utf-8",
                )
                converted += 1
            except Exception as error:  # noqa: BLE001 - 单个文件失败不中断整体
                failed.append(f"{public_path}: {error}")

    # Scratch files are never publication artifacts. Remove converted legacy
    # decks before the release tree is validated and made immutable.
    shutil.rmtree(scratch, ignore_errors=True)

    print(
        f"convert-henukit-slides: converted {converted}, skipped {skipped}, "
        f"failed {len(failed)}",
        file=sys.stderr,
    )
    for message in failed[:20]:
        print(f"convert-henukit-slides: FAILED {message}", file=sys.stderr)
    if len(failed) > 20:
        print(f"convert-henukit-slides: ... and {len(failed) - 20} more", file=sys.stderr)
    if failed:
        sys.exit(5)


if __name__ == "__main__":
    main()
